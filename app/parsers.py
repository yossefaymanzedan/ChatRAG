import re
from dataclasses import dataclass, field
from pathlib import Path
from typing import Any

from docx import Document as DocxDocument
from pptx import Presentation
from pypdf import PdfReader

try:
    import numpy as np
except Exception:
    np = None

try:
    import pymupdf
except Exception:
    try:
        import fitz as pymupdf  # Backward-compatible PyMuPDF import.
    except Exception:
        pymupdf = None

try:
    from rapidocr import RapidOCR
except Exception:
    RapidOCR = None


SUPPORTED_EXTS = {".pdf", ".docx", ".pptx", ".txt", ".md"}
OCR_STAT_KEYS = (
    "pdf_files_scanned",
    "images_found",
    "image_pages_found",
    "image_pages_processed",
    "image_pages_added_to_rag",
    "image_pages_not_added_to_rag",
)


@dataclass
class ParseResult:
    chunks: list[dict]
    warnings: list[str]
    stats: dict[str, int] = field(default_factory=dict)


_OCR_ENGINE: Any | None = None
_OCR_ENGINE_INITIALIZED = False


def _empty_ocr_stats() -> dict[str, int]:
    return {key: 0 for key in OCR_STAT_KEYS}


def _format_page_list(pages: list[int], limit: int = 12) -> str:
    if not pages:
        return "-"
    unique = sorted(set(int(p) for p in pages if isinstance(p, int)))
    if len(unique) <= limit:
        return ", ".join(str(p) for p in unique)
    head = ", ".join(str(p) for p in unique[:limit])
    return f"{head} (+{len(unique) - limit} more)"


def _get_ocr_engine() -> tuple[Any | None, str | None]:
    global _OCR_ENGINE, _OCR_ENGINE_INITIALIZED
    if _OCR_ENGINE_INITIALIZED:
        if _OCR_ENGINE is None:
            return None, "RapidOCR engine initialization failed"
        return _OCR_ENGINE, None
    _OCR_ENGINE_INITIALIZED = True
    if RapidOCR is None:
        return None, "rapidocr is not installed"
    try:
        _OCR_ENGINE = RapidOCR()
        return _OCR_ENGINE, None
    except Exception as exc:
        _OCR_ENGINE = None
        return None, f"RapidOCR init error: {exc}"


def _collect_ocr_text(node: Any, parts: list[str]) -> None:
    if node is None:
        return
    if isinstance(node, str):
        text = node.strip()
        if text:
            parts.append(text)
        return
    for attr in ("txts", "texts", "text"):
        if hasattr(node, attr):
            try:
                attr_value = getattr(node, attr)
            except Exception:
                attr_value = None
            _collect_ocr_text(attr_value, parts)
            if parts:
                return
    if isinstance(node, dict):
        for key in ("text", "txt", "transcription"):
            value = node.get(key)
            if isinstance(value, str) and value.strip():
                parts.append(value.strip())
                return
        for value in node.values():
            _collect_ocr_text(value, parts)
        return
    if isinstance(node, (list, tuple)):
        if len(node) >= 2 and isinstance(node[1], str) and not isinstance(node[0], str):
            text = node[1].strip()
            if text:
                parts.append(text)
            return
        for item in node:
            _collect_ocr_text(item, parts)


def _extract_ocr_text(raw_output: Any) -> str:
    payload = raw_output
    if isinstance(raw_output, tuple) and raw_output:
        payload = raw_output[0]
    parts: list[str] = []
    _collect_ocr_text(payload, parts)
    return "\n".join(parts).strip()


def _merge_extracted_and_ocr_text(extracted_text: str, ocr_text: str) -> tuple[str, bool]:
    base = (extracted_text or "").strip()
    ocr = (ocr_text or "").strip()
    if not ocr:
        return base, False
    if not base:
        return ocr, True

    merged_lines = [ln.strip() for ln in base.splitlines() if ln.strip()]
    seen = {
        re.sub(r"\s+", " ", ln).strip().lower()
        for ln in merged_lines
        if ln.strip()
    }
    added = False
    for ln in [x.strip() for x in ocr.splitlines() if x.strip()]:
        sig = re.sub(r"\s+", " ", ln).strip().lower()
        if not sig or sig in seen:
            continue
        merged_lines.append(ln)
        seen.add(sig)
        added = True

    if not added:
        return base, False
    return "\n".join(merged_lines), True


def _open_pdf_image_context(path: Path) -> tuple[Any | None, dict[int, int], str | None]:
    if pymupdf is None:
        return None, {}, "PyMuPDF is not installed"
    try:
        doc = pymupdf.open(str(path))
    except Exception as exc:
        return None, {}, f"PyMuPDF open failed: {exc}"
    image_counts: dict[int, int] = {}
    for idx in range(doc.page_count):
        page_num = idx + 1
        try:
            page = doc.load_page(idx)
            images = page.get_images(full=True) or []
            if images:
                image_counts[page_num] = len(images)
        except Exception:
            continue
    return doc, image_counts, None


def _ocr_pdf_page(doc: Any, page_index: int, render_zoom: float) -> str:
    if np is None:
        raise RuntimeError("numpy is not available for OCR image conversion")
    ocr_engine, ocr_err = _get_ocr_engine()
    if ocr_engine is None:
        raise RuntimeError(ocr_err or "RapidOCR is unavailable")
    page = doc.load_page(page_index)
    zoom = max(1.0, float(render_zoom))
    matrix = pymupdf.Matrix(zoom, zoom)
    pix = page.get_pixmap(matrix=matrix, alpha=False)
    image = np.frombuffer(pix.samples, dtype=np.uint8).reshape(pix.height, pix.width, pix.n)
    if image.ndim == 3 and image.shape[2] == 4:
        image = image[:, :, :3]
    result = ocr_engine(image)
    return _extract_ocr_text(result)


def token_count(text: str) -> int:
    return len(text.split())


def preview_text(text: str, max_len: int = 280) -> str:
    text = re.sub(r"\s+", " ", text).strip()
    if len(text) <= max_len:
        return text
    return text[: max_len - 3].rstrip() + "..."


def split_blocks(text: str) -> list[str]:
    raw = [b.strip() for b in re.split(r"\n\s*\n+", text) if b.strip()]
    if raw:
        return raw
    line_blocks = [line.strip() for line in text.splitlines() if line.strip()]
    return line_blocks


def _looks_like_toc_lines(lines: list[str]) -> bool:
    toc_like_lines = 0
    for ln in lines[:24]:
        has_dot_leader = ("..." in ln) or (" . " in ln)
        ends_with_page = bool(re.search(r"\b\d{1,4}\s*$", ln))
        starts_with_index = bool(re.match(r"^\s*(chapter|part|appendix|\d+(\.\d+){0,3})\b", ln.lower()))
        if ends_with_page and (has_dot_leader or starts_with_index):
            toc_like_lines += 1
    return toc_like_lines >= 3


def _looks_like_front_matter_heading(text: str) -> bool:
    heading = re.sub(r"\s+", " ", (text or "").strip().lower())
    if not heading:
        return False
    prefixes = (
        "acknowledgment",
        "acknowledgement",
        "preface",
        "foreword",
        "copyright",
        "about the author",
        "about this book",
        "list of figures",
        "list of tables",
        "table of contents",
        "contents",
    )
    return any(heading.startswith(p) for p in prefixes)


def _looks_like_acknowledgment_list(text: str) -> bool:
    lines = [ln.strip() for ln in (text or "").splitlines() if ln.strip()]
    if len(lines) < 8:
        return False
    bullet_lines = [ln for ln in lines if ln.startswith(("•", "-", "*"))]
    if len(bullet_lines) < 6:
        return False
    low = re.sub(r"\s+", " ", (text or "").lower())
    signal_terms = ("typo", "suggest", "correction", "chapter", "thanks", "thank")
    signal_hits = sum(1 for t in signal_terms if t in low)
    return signal_hits >= 2


def looks_like_toc_block(text: str) -> bool:
    low = re.sub(r"\s+", " ", (text or "").strip().lower())
    if not low:
        return False

    if "table of contents" in low or low.startswith("contents"):
        return True

    lines = [ln.strip() for ln in (text or "").splitlines() if ln.strip()]
    if len(lines) < 3:
        return False

    return _looks_like_toc_lines(lines)


def _detect_main_content_start(page_texts: list[tuple[int, str]], scan_pages: int) -> int | None:
    for page_num, text in page_texts[: max(1, scan_pages)]:
        lines = [ln.strip() for ln in text.splitlines() if ln.strip()]
        head = " ".join(lines[:5]).lower()
        words = token_count(text)
        chapter_like = bool(
            re.search(r"\b(chapter\s+\d+|part\s+[ivx0-9]+|introduction)\b", head)
        )
        if chapter_like and words >= 80 and not _looks_like_toc_lines(lines):
            return page_num
    return None


def _is_front_matter_block(block: str, page_num: int, start_page: int | None, scan_pages: int) -> bool:
    lines = [ln.strip() for ln in block.splitlines() if ln.strip()]
    if not lines:
        return False
    low = re.sub(r"\s+", " ", block.strip().lower())

    if looks_like_toc_block(block):
        return True

    front_window = max(1, scan_pages)
    if start_page is not None:
        front_window = min(front_window, max(1, start_page - 1))
    in_front_range = page_num <= front_window
    short_block = token_count(block) <= 220

    if in_front_range and _looks_like_front_matter_heading(low):
        return True

    if in_front_range and short_block and re.match(r"^(isbn|copyright|all rights reserved)\b", low):
        return True

    if in_front_range and _looks_like_acknowledgment_list(block):
        return True

    return False


def _is_front_matter_page(text: str, page_num: int, start_page: int | None, scan_pages: int) -> bool:
    front_window = max(1, scan_pages)
    if start_page is not None:
        front_window = min(front_window, max(1, start_page - 1))
    if page_num > front_window:
        return False

    lines = [ln.strip() for ln in (text or "").splitlines() if ln.strip()]
    if not lines:
        return False
    head = " ".join(lines[:8])
    if _looks_like_front_matter_heading(head):
        return True
    if _looks_like_toc_lines(lines):
        return True
    if _looks_like_acknowledgment_list(text):
        return True
    return False


def parse_file(
    path: Path,
    pdf_text_threshold: int = 40,
    ignore_front_matter: bool = True,
    front_matter_scan_pages: int = 24,
    enable_pdf_ocr: bool = True,
    pdf_ocr_render_zoom: float = 2.0,
) -> ParseResult:
    ext = path.suffix.lower()
    if ext == ".pdf":
        return parse_pdf(
            path,
            pdf_text_threshold,
            ignore_front_matter=ignore_front_matter,
            front_matter_scan_pages=front_matter_scan_pages,
            enable_ocr=enable_pdf_ocr,
            ocr_render_zoom=pdf_ocr_render_zoom,
        )
    if ext == ".docx":
        return parse_docx(path)
    if ext == ".pptx":
        return parse_pptx(path)
    if ext == ".md":
        return parse_markdown(path)
    if ext == ".txt":
        return parse_txt(path)
    return ParseResult(chunks=[], warnings=[f"Unsupported file extension: {ext}"])


def parse_pdf(
    path: Path,
    threshold: int,
    ignore_front_matter: bool = True,
    front_matter_scan_pages: int = 24,
    enable_ocr: bool = True,
    ocr_render_zoom: float = 2.0,
) -> ParseResult:
    chunks: list[dict] = []
    warnings: list[str] = []
    ocr_stats = _empty_ocr_stats()
    ocr_stats["pdf_files_scanned"] = 1
    idx = 0
    skipped_toc_blocks = 0
    skipped_front_matter_blocks = 0
    skipped_front_matter_pages = 0
    ocr_added_pages: list[int] = []
    ocr_failed_pages: list[int] = []
    ocr_no_text_pages: list[int] = []
    low_text_skipped_pages: list[int] = []
    page_counts_with_chunks: dict[int, int] = {}

    try:
        reader = PdfReader(str(path))
    except Exception as exc:
        return ParseResult(
            chunks=[],
            warnings=[f"{path.name}: failed to open PDF ({exc})"],
            stats=ocr_stats,
        )

    pdf_image_doc, image_counts, image_scan_warning = _open_pdf_image_context(path)
    image_pages = {page for page, count in image_counts.items() if int(count or 0) > 0}
    ocr_stats["image_pages_found"] = len(image_pages)
    ocr_stats["images_found"] = sum(int(c or 0) for c in image_counts.values())
    if image_scan_warning:
        warnings.append(f"{path.name}: image scan unavailable ({image_scan_warning})")

    ocr_unavailable_reason: str | None = None
    if enable_ocr:
        if pdf_image_doc is None:
            ocr_unavailable_reason = "PyMuPDF context is unavailable"
        elif np is None:
            ocr_unavailable_reason = "numpy is unavailable"
        else:
            _ocr_engine, ocr_err = _get_ocr_engine()
            if _ocr_engine is None:
                ocr_unavailable_reason = ocr_err or "RapidOCR is unavailable"

    if enable_ocr and image_pages and ocr_unavailable_reason:
        warnings.append(f"{path.name}: OCR unavailable ({ocr_unavailable_reason})")

    ocr_accept_threshold = max(20, int(threshold / 2))
    start_page = None
    page_texts: list[tuple[int, str]] = []
    try:
        for page_num, page in enumerate(reader.pages, start=1):
            has_images = page_num in image_pages
            try:
                text = page.extract_text() or ""
            except Exception as exc:
                warnings.append(f"{path.name}: page {page_num} parse error ({exc})")
                continue
            text = text.strip()
            ocr_text = ""
            ocr_attempted = False
            if (
                enable_ocr
                and has_images
                and not ocr_unavailable_reason
                and pdf_image_doc is not None
            ):
                ocr_attempted = True
                ocr_stats["image_pages_processed"] += 1
                try:
                    ocr_text = _ocr_pdf_page(
                        pdf_image_doc,
                        page_index=page_num - 1,
                        render_zoom=ocr_render_zoom,
                    ).strip()
                except Exception:
                    ocr_failed_pages.append(page_num)
                    ocr_text = ""

            if len(text) < threshold:
                if ocr_attempted:
                    if len(ocr_text) >= ocr_accept_threshold:
                        text = ocr_text
                        ocr_added_pages.append(page_num)
                    else:
                        ocr_no_text_pages.append(page_num)
                        continue
                else:
                    low_text_skipped_pages.append(page_num)
                    continue
            else:
                if ocr_attempted and ocr_text:
                    merged_text, added_from_ocr = _merge_extracted_and_ocr_text(text, ocr_text)
                    text = merged_text
                    if added_from_ocr:
                        ocr_added_pages.append(page_num)
            page_texts.append((page_num, text))

        if ignore_front_matter:
            start_page = _detect_main_content_start(page_texts, front_matter_scan_pages)

        ocr_added_page_set = set(ocr_added_pages)
        for page_num, text in page_texts:
            is_image_source_page = page_num in ocr_added_page_set
            if ignore_front_matter and _is_front_matter_page(
                text,
                page_num=page_num,
                start_page=start_page,
                scan_pages=front_matter_scan_pages,
            ):
                skipped_front_matter_pages += 1
                continue
            blocks = split_blocks(text)
            chunked_on_page = 0
            for para_idx, block in enumerate(blocks, start=1):
                if looks_like_toc_block(block):
                    skipped_toc_blocks += 1
                    continue
                if ignore_front_matter and _is_front_matter_block(
                    block,
                    page_num=page_num,
                    start_page=start_page,
                    scan_pages=front_matter_scan_pages,
                ):
                    skipped_front_matter_blocks += 1
                    continue
                chunks.append(
                    {
                        "chunk_index": idx,
                        "text": block,
                        "token_count": token_count(block),
                        "anchor_type": "image" if is_image_source_page else "pdf_page",
                        "anchor_page": page_num,
                        "anchor_section": "Image source" if is_image_source_page else None,
                        "anchor_paragraph": para_idx,
                        "anchor_row": None,
                        "start_char": None,
                        "end_char": None,
                        "preview": preview_text(block),
                    }
                )
                idx += 1
                chunked_on_page += 1
            if chunked_on_page:
                page_counts_with_chunks[page_num] = chunked_on_page
    finally:
        if pdf_image_doc is not None:
            try:
                pdf_image_doc.close()
            except Exception:
                pass

    if skipped_toc_blocks:
        warnings.append(f"{path.name}: skipped {skipped_toc_blocks} TOC-like block(s)")
    if skipped_front_matter_blocks:
        warnings.append(f"{path.name}: skipped {skipped_front_matter_blocks} front-matter block(s)")
    if skipped_front_matter_pages:
        warnings.append(f"{path.name}: skipped {skipped_front_matter_pages} front-matter page(s)")
    if ignore_front_matter and start_page is not None:
        warnings.append(f"{path.name}: detected main content start around page {start_page}")

    if low_text_skipped_pages:
        warnings.append(
            f"{path.name}: skipped {len(low_text_skipped_pages)} low-text page(s)"
            f" (pages: {_format_page_list(low_text_skipped_pages)})"
        )
    if ocr_added_pages:
        warnings.append(
            f"{path.name}: OCR extracted usable text from {len(set(ocr_added_pages))} page(s)"
            f" (pages: {_format_page_list(ocr_added_pages)})"
        )
    if ocr_failed_pages:
        warnings.append(
            f"{path.name}: OCR failed on {len(set(ocr_failed_pages))} page(s)"
            f" (pages: {_format_page_list(ocr_failed_pages)})"
        )
    if ocr_no_text_pages:
        warnings.append(
            f"{path.name}: OCR found no usable text on {len(set(ocr_no_text_pages))} page(s)"
            f" (pages: {_format_page_list(ocr_no_text_pages)})"
        )

    added_image_pages = len(image_pages.intersection(page_counts_with_chunks.keys()))
    not_added_image_pages = max(0, len(image_pages) - added_image_pages)
    ocr_stats["image_pages_added_to_rag"] = added_image_pages
    ocr_stats["image_pages_not_added_to_rag"] = not_added_image_pages

    if image_pages:
        warnings.append(
            (
                f"{path.name}: image/OCR summary -> images found: {ocr_stats['images_found']} | "
                f"image pages: {ocr_stats['image_pages_found']} | OCR processed pages: {ocr_stats['image_pages_processed']} | "
                f"image pages added to RAG: {ocr_stats['image_pages_added_to_rag']} | "
                f"image pages not added: {ocr_stats['image_pages_not_added_to_rag']}"
            )
        )

    return ParseResult(chunks=chunks, warnings=warnings, stats=ocr_stats)


def parse_docx(path: Path) -> ParseResult:
    doc = DocxDocument(str(path))
    para_records: list[tuple[int, str, str | None]] = []
    nearest_heading: str | None = None

    for i, para in enumerate(doc.paragraphs):
        text = para.text.strip()
        if not text:
            continue
        style = (para.style.name or "").lower() if para.style else ""
        if style.startswith("heading"):
            nearest_heading = text
            continue
        para_records.append((i, text, nearest_heading))

    merged: list[tuple[int, str, str | None]] = []
    i = 0
    while i < len(para_records):
        start_idx, text, section = para_records[i]
        while len(text) < 80 and i + 1 < len(para_records):
            i += 1
            _, nxt_text, _ = para_records[i]
            text = f"{text}\n{nxt_text}"
            if len(text) >= 160:
                break
        merged.append((start_idx, text, section))
        i += 1

    chunks: list[dict] = []
    for chunk_idx, (paragraph_index, text, section) in enumerate(merged):
        chunks.append(
            {
                "chunk_index": chunk_idx,
                "text": text,
                "token_count": token_count(text),
                "anchor_type": "docx_paragraph",
                "anchor_page": None,
                "anchor_section": section,
                "anchor_paragraph": paragraph_index,
                "anchor_row": None,
                "start_char": None,
                "end_char": None,
                "preview": preview_text(text),
            }
        )

    return ParseResult(chunks=chunks, warnings=[])


def parse_pptx(path: Path) -> ParseResult:
    prs = Presentation(str(path))
    chunks: list[dict] = []
    chunk_idx = 0
    for slide_num, slide in enumerate(prs.slides, start=1):
        block_idx = 1
        for shape in slide.shapes:
            text = getattr(shape, "text", "") or ""
            text = text.strip()
            if not text:
                continue
            chunks.append(
                {
                    "chunk_index": chunk_idx,
                    "text": text,
                    "token_count": token_count(text),
                    "anchor_type": "ppt_slide",
                    "anchor_page": slide_num,
                    "anchor_section": f"Slide {slide_num}",
                    "anchor_paragraph": block_idx,
                    "anchor_row": None,
                    "start_char": None,
                    "end_char": None,
                    "preview": preview_text(text),
                }
            )
            chunk_idx += 1
            block_idx += 1
    return ParseResult(chunks=chunks, warnings=[])


def parse_markdown(path: Path) -> ParseResult:
    text = path.read_text(encoding="utf-8", errors="ignore")
    lines = text.splitlines()
    chunks: list[dict] = []
    current_heading = "Document"
    block: list[str] = []
    block_index = 0

    def flush() -> None:
        nonlocal block, block_index
        content = "\n".join(block).strip()
        if not content:
            block = []
            return
        chunks.append(
            {
                "chunk_index": len(chunks),
                "text": content,
                "token_count": token_count(content),
                "anchor_type": "md_heading",
                "anchor_page": None,
                "anchor_section": current_heading,
                "anchor_paragraph": block_index,
                "anchor_row": None,
                "start_char": None,
                "end_char": None,
                "preview": preview_text(content),
            }
        )
        block_index += 1
        block = []

    for line in lines:
        heading = re.match(r"^(#{1,6})\s+(.+)$", line.strip())
        if heading:
            flush()
            current_heading = heading.group(2).strip()
            block_index = 0
            continue
        if line.strip() == "":
            flush()
        else:
            block.append(line)
    flush()

    return ParseResult(chunks=chunks, warnings=[])


def parse_txt(path: Path) -> ParseResult:
    text = path.read_text(encoding="utf-8", errors="ignore")
    blocks = split_blocks(text)
    chunks = []
    for idx, block in enumerate(blocks):
        chunks.append(
            {
                "chunk_index": idx,
                "text": block,
                "token_count": token_count(block),
                "anchor_type": "txt_block",
                "anchor_page": None,
                "anchor_section": None,
                "anchor_paragraph": idx,
                "anchor_row": None,
                "start_char": None,
                "end_char": None,
                "preview": preview_text(block),
            }
        )

    return ParseResult(chunks=chunks, warnings=[])
